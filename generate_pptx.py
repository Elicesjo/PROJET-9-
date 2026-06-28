from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

BLUE_DARK  = RGBColor(0x1A, 0x52, 0x76)
BLUE_MED   = RGBColor(0x2E, 0x86, 0xC1)
BLUE_LIGHT = RGBColor(0xD6, 0xEA, 0xF8)
GRAY_BG    = RGBColor(0xF2, 0xF2, 0xF2)
GRAY_TEXT  = RGBColor(0x88, 0x88, 0x88)
GREEN      = RGBColor(0x1E, 0x8A, 0x44)
RED        = RGBColor(0xC0, 0x39, 0x2B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1C, 0x1C, 0x1C)

FONT = "Aptos Display"

SLIDE_W       = 33.87
SLIDE_H       = 19.05
MARGIN        = 2.0
CONTENT_W     = SLIDE_W - 2 * MARGIN   # 29.87
LEFT_COL_W    = CONTENT_W * 3 / 5      # 17.92
GAP           = 0.4
RIGHT_COL_L   = MARGIN + LEFT_COL_W + GAP
BOX_W         = CONTENT_W * 2 / 5 - 1.8

TITLE_LEFT    = MARGIN
TITLE_TOP     = 0.7
TITLE_H_1LINE = 1.8    # height for a single-line title
TITLE_H_2LINE = 3.2    # height for a two-line title
LISERET_GAP   = 1.5    # gap above and below the liseret
LISERET_H     = 0.08
BAND_H        = 2.5    # card header band height (fits 1–2 lines at 16 pt)

BASE       = os.path.dirname(os.path.abspath(__file__))
MOCKUP_DIR = os.path.join(BASE, "mockup")


# ── Primitives ────────────────────────────────────────────────────────────────

def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _no_border(shape):
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    for child in list(ln):
        ln.remove(child)
    etree.SubElement(ln, qn("a:noFill"))


def _rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, Cm(left), Cm(top), Cm(width), Cm(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    _no_border(shape)
    return shape


def _oval(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(9, Cm(left), Cm(top), Cm(width), Cm(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    _no_border(shape)
    return shape


def _set_run(run, text, size, bold, color):
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _textbox(slide, text, left, top, width, height,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _set_run(p.add_run(), text, size, bold, color)
    return txb


def _multiline(slide, lines, left, top, width, height, size=18, space_after=5):
    txb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, (text, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        _set_run(p.add_run(), text, size, bold, color if color else DARK)
        if text.strip():
            p.space_after = Pt(space_after)
    return txb


# ── Composants ────────────────────────────────────────────────────────────────

def add_title(slide, text):
    """Draw title + liseret. Returns (content_top, content_h) computed dynamically."""
    two_lines = len(text) > 44
    title_h    = TITLE_H_2LINE if two_lines else TITLE_H_1LINE
    liseret_top = TITLE_TOP + title_h + LISERET_GAP
    content_top = liseret_top + LISERET_H + LISERET_GAP
    content_h   = SLIDE_H - content_top - 1.0

    txb = slide.shapes.add_textbox(
        Cm(TITLE_LEFT), Cm(TITLE_TOP), Cm(CONTENT_W), Cm(title_h)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_run(p.add_run(), text, 36, True, DARK)

    _rect(slide, TITLE_LEFT, liseret_top, CONTENT_W, LISERET_H, BLUE_DARK)
    return content_top, content_h


def add_page_number(slide, number, is_cover=False):
    color = WHITE if is_cover else DARK
    _textbox(slide, str(number),
             SLIDE_W - 2.2, SLIDE_H - 1.1, 1.8, 0.8,
             size=14, color=color, align=PP_ALIGN.RIGHT)


def add_kpi_box(slide, value, label, left, top, width, height, value_color=BLUE_MED):
    _rect(slide, left, top, width, height, GRAY_BG)
    pad = 0.4
    _textbox(slide, value,
             left + pad, top + height * 0.1, width - 2 * pad, height * 0.5,
             size=40, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    _textbox(slide, label,
             left + pad, top + height * 0.62, width - 2 * pad, height * 0.33,
             size=14, color=GRAY_TEXT, align=PP_ALIGN.CENTER)


def add_card(slide, title, body_lines, left, top, width, height):
    _rect(slide, left, top, width, height, GRAY_BG)
    _rect(slide, left, top, width, BAND_H, BLUE_DARK)
    pad = 0.4
    # Vertically center title within band
    title_top = top + (BAND_H - 1.6) / 2
    txb = slide.shapes.add_textbox(
        Cm(left + pad), Cm(title_top), Cm(width - 2 * pad), Cm(1.6)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_run(p.add_run(), title, 16, True, WHITE)

    body_top = top + BAND_H + 0.4
    body_h   = height - BAND_H - 0.6
    _multiline(slide, body_lines,
               left + pad, body_top, width - 2 * pad, body_h,
               size=16, space_after=6)


def add_pipeline_step(slide, number, title, body, left, top, width, height):
    _rect(slide, left, top, width, height, BLUE_LIGHT)
    circ_size = 1.1
    circ_top  = top + (height - circ_size) / 2
    circ = _oval(slide, left + 0.4, circ_top, circ_size, circ_size, BLUE_DARK)
    tf = circ.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), str(number), 18, True, WHITE)
    text_left = left + 0.4 + circ_size + 0.4
    text_w    = width - 0.4 - circ_size - 0.8
    _textbox(slide, title,
             text_left, top + 0.25, text_w, height * 0.42,
             size=18, bold=True, color=DARK)
    _textbox(slide, body,
             text_left, top + height * 0.48, text_w, height * 0.48,
             size=16, color=GRAY_TEXT)


def add_section_divider(prs, title, subtitle, page_num):
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, BLUE_DARK)
    _textbox(slide, title,
             3.0, 6.5, SLIDE_W - 6.0, 3.5,
             size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        _textbox(slide, subtitle,
                 3.0, 10.5, SLIDE_W - 6.0, 2.0,
                 size=24, color=RGBColor(0x5D, 0xAD, 0xE8), align=PP_ALIGN.CENTER)
    add_page_number(slide, page_num, is_cover=True)


def add_file_slide(prs, title, file_name, analyse, modifs, page_num):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, title)

    left_lines = (
        [("Analyse du fichier :", True, BLUE_DARK), ("", False, None)]
        + [(f"• {a}", False, None) for a in analyse]
        + [("", False, None),
           ("Modifications effectuées :", True, BLUE_DARK),
           ("", False, None)]
        + [(f"• {m}", False, None) for m in modifs]
    )
    _multiline(slide, left_lines, MARGIN, ct, LEFT_COL_W, ch, size=18, space_after=6)

    box_h = ch * 0.72
    _rect(slide, RIGHT_COL_L, ct, BOX_W, box_h, GRAY_BG)
    pad = 0.4
    _textbox(slide, file_name,
             RIGHT_COL_L + pad, ct + 0.4, BOX_W - 2 * pad, 1.1,
             size=14, bold=True, color=BLUE_DARK)
    for i, a in enumerate(analyse):
        _textbox(slide, f"• {a}",
                 RIGHT_COL_L + pad, ct + 1.7 + i * 1.1, BOX_W - 2 * pad, 1.0,
                 size=14, color=GRAY_TEXT)

    add_page_number(slide, page_num)


# ── Slides ────────────────────────────────────────────────────────────────────

def build_cover(prs):
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, BLUE_DARK)
    _textbox(slide,
             "Accès à l'eau potable dans le monde",
             3.0, 5.5, SLIDE_W - 6.0, 4.0,
             size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide,
             "Projet DWFA  ·  Analyse exploratoire et tableau de bord interactif",
             3.0, 10.2, SLIDE_W - 6.0, 2.5,
             size=36, color=WHITE, align=PP_ALIGN.CENTER)
    _rect(slide, MARGIN, 14.5, CONTENT_W, 0.06, RGBColor(0x5D, 0xAD, 0xE8))
    add_page_number(slide, 1, is_cover=True)


def build_objectif(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "Objectif de la mission")

    intro_lines = [
        ("DWFA — Drinking Water For All — œuvre pour améliorer l'accès à l'eau "
         "potable dans les pays en développement.", False, None),
        ("", False, None),
        ("La mission : identifier le pays le plus en difficulté dans chaque "
         "domaine, afin d'orienter l'intervention.", False, None),
    ]
    _multiline(slide, intro_lines, MARGIN, ct, CONTENT_W, 3.8, size=18, space_after=5)

    domains = [
        ("D1 — Construction",
         [("Accès basique < 50 %", False, None),
          ("Pas d'infrastructure de base", False, None),
          ("Besoin : créer le réseau", False, None),
          ("Critère : population non couverte", False, None)]),
        ("D2 — Modernisation",
         [("Accès basique 50–90 %", False, None),
          ("Infrastructure partielle", False, None),
          ("Besoin : étendre et sécuriser", False, None),
          ("Critère : écart basique / sécurisé", False, None)]),
        ("D3 — Conseil gouvernemental",
         [("Instabilité politique élevée", False, None),
          ("Cadre institutionnel défaillant", False, None),
          ("Besoin : accompagnement politique", False, None),
          ("Critère : score WGI stabilité", False, None)]),
    ]

    n = len(domains)
    card_gap = 0.4
    card_w = (CONTENT_W - (n - 1) * card_gap) / n
    card_top = ct + 4.2
    card_h = ch - 4.2

    for i, (title, body) in enumerate(domains):
        add_card(slide, title, body,
                 MARGIN + i * (card_w + card_gap), card_top, card_w, card_h)

    add_page_number(slide, 2)


def build_mission(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "Mission — 3 axes d'analyse")

    steps = [
        ("Exploration des données",
         [("Structure et qualité des 5 fichiers sources", False, None),
          ("Types, valeurs manquantes, anomalies", False, None),
          ("Couverture géographique et temporelle", False, None)]),
        ("Création des dashboards",
         [("Vue Mondiale : KPIs et cartes", False, None),
          ("Vue Continentale : comparaison régionale", False, None),
          ("Vue Nationale : zoom pays, évolution", False, None)]),
        ("Choix des pays prioritaires",
         [("Sélection par domaine D1 / D2 / D3", False, None),
          ("Critères basés sur les indicateurs", False, None),
          ("Recommandations actionnables DWFA", False, None)]),
    ]

    n = len(steps)
    card_gap = 0.4
    card_w = (CONTENT_W - (n - 1) * card_gap) / n

    for i, (title, body) in enumerate(steps):
        add_card(slide, title, body,
                 MARGIN + i * (card_w + card_gap), ct, card_w, ch)

    add_page_number(slide, 3)


def build_outils(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "Outils & Ressources")

    outils_lines = [
        ("Environnement d'analyse", True, BLUE_DARK),
        ("", False, None),
        ("• Python 3.12 — pandas, numpy, python-pptx", False, None),
        ("• Jupyter Notebook — exploration interactive", False, None),
        ("• uv — gestion des dépendances", False, None),
        ("", False, None),
        ("Visualisation & Dashboard", True, BLUE_DARK),
        ("", False, None),
        ("• Tableau Public — dashboards interactifs", False, None),
        ("• 3 vues interconnectées : Mondiale / Continentale / Nationale", False, None),
        ("• Publication sans infrastructure serveur", False, None),
        ("", False, None),
        ("Données sources", True, BLUE_DARK),
        ("", False, None),
        ("• 5 fichiers OMS / Banque Mondiale", False, None),
        ("• 194 pays · 6 régions OMS · 2000–2017", False, None),
    ]
    _multiline(slide, outils_lines, MARGIN, ct, LEFT_COL_W, ch, size=18, space_after=5)

    kpi_h = (ch - 2 * 0.4) / 3
    kpis = [
        ("5",      "Fichiers sources OMS / BM",  BLUE_MED),
        ("194",    "Pays couverts",               BLUE_MED),
        ("18 ans", "Période 2000–2017",           BLUE_MED),
    ]
    for i, (val, lbl, col) in enumerate(kpis):
        add_kpi_box(slide, val, lbl,
                    RIGHT_COL_L, ct + i * (kpi_h + 0.4), BOX_W, kpi_h, col)

    add_page_number(slide, 3)


def build_etapes(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "Grandes étapes du projet")

    etapes = [
        ("Nettoyage\n& Fusion",
         "Exploration, types, doublons, jointures des 5 sources"),
        ("Blueprint\n& Mock-Up",
         "Conception des dashboards : structure, vues, indicateurs"),
        ("Enrichissement\ndes données",
         "Indicateurs dérivés : taux pop. rurale, décès WASH"),
        ("Construction\nvues globales",
         "Dashboards Mondiale, Continentale, Nationale"),
        ("Construction\nvues spécifiques",
         "Analyses D1, D2, D3 — sélection pays prioritaires"),
    ]

    n = len(etapes)
    arrow_w  = 0.65
    card_gap = 0.3
    card_w   = (CONTENT_W - (n - 1) * (arrow_w + card_gap)) / n

    for i, (step_title, step_body) in enumerate(etapes):
        left = MARGIN + i * (card_w + card_gap + arrow_w)
        _rect(slide, left, ct, card_w, ch, BLUE_LIGHT)
        circ = _oval(slide, left + card_w / 2 - 0.65, ct + 0.45, 1.3, 1.3, BLUE_DARK)
        tf = circ.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), str(i + 1), 18, True, WHITE)
        _textbox(slide, step_title,
                 left + 0.3, ct + 2.1, card_w - 0.6, 2.4,
                 size=16, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
        _textbox(slide, step_body,
                 left + 0.3, ct + 4.8, card_w - 0.6, ch - 5.0,
                 size=14, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
        if i < n - 1:
            _textbox(slide, "→",
                     left + card_w + card_gap / 2, ct + ch / 2 - 0.5, arrow_w, 1.0,
                     size=22, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

    add_page_number(slide, 4)


def build_exploration_fichiers(prs, page_num):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "Sources — 5 fichiers OMS et Banque Mondiale")

    fichiers = [
        ("Population.csv",
         "194 pays · Urban / Rural · 2000–2017",
         "Multiply ×1 000, typage float64, 0 doublon"),
        ("MortalityRateAttributedToWater.csv",
         "183 pays · Total · 2016 uniquement",
         "Typage, doublons supprimés, jointure sur Pays"),
        ("BasicAndSafelyManagedDrinkingWaterServices.csv",
         "194 pays · Total / Urban / Rural · 2000–2017",
         "Séparation granularité, 50 % NaN sécurisé (manque OMS)"),
        ("PoliticalStability.csv",
         "190 pays · Score WGI −3,31 à +1,76 · 2000–2017",
         "Typage, doublons, années partielles documentées"),
        ("RegionCountry.csv",
         "194 pays · 6 régions OMS",
         "Clé jointure validée, 0 doublon"),
    ]

    n = len(fichiers)
    row_gap = 0.35
    row_h = (ch - (n - 1) * row_gap) / n
    accent_w = 0.45
    pad = 0.5
    col1_w = 12.0
    col2_w = 9.0
    col1_left = MARGIN + accent_w + pad
    col2_left = col1_left + col1_w + 0.4
    col3_left = col2_left + col2_w + 0.4
    col3_w = SLIDE_W - MARGIN - col3_left

    for i, (fname, stats, transforms) in enumerate(fichiers):
        top = ct + i * (row_h + row_gap)
        _rect(slide, MARGIN, top, CONTENT_W, row_h, GRAY_BG)
        _rect(slide, MARGIN, top, accent_w, row_h, BLUE_DARK)
        text_top = top + 0.25
        _textbox(slide, fname, col1_left, text_top, col1_w, row_h * 0.65,
                 size=14, bold=True, color=BLUE_DARK)
        _textbox(slide, stats, col2_left, text_top, col2_w, row_h * 0.65,
                 size=13, color=DARK)
        _textbox(slide, f"→  {transforms}", col3_left, text_top, col3_w, row_h * 0.65,
                 size=13, color=GRAY_TEXT)

    add_page_number(slide, page_num)


def build_mockup_globales(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "Mock-Up — Analyses globales")

    vues = [
        ("Vue Mondiale",
         [("3 KPIs (accès basique, sécurisé, décès)", False, None),
          ("2 filtres : région, année", False, None),
          ("Graphiques :", True, BLUE_DARK),
          ("• Carte mondiale mortalité WASH", False, None),
          ("• Bar instabilité par région", False, None),
          ("• Line chart évolution accès", False, None)]),
        ("Vue Continentale",
         [("KPIs filtrés par région", False, None),
          ("2 filtres : région, pays", False, None),
          ("Graphiques :", True, BLUE_DARK),
          ("• Scatter mortalité vs accès", False, None),
          ("• Scatter stabilité vs sécurisé", False, None),
          ("• Bar classement pays", False, None),
          ("• Line évolution par pays", False, None)]),
        ("Vue Nationale",
         [("KPIs filtrés par pays", False, None),
          ("2 filtres : pays, indicateur", False, None),
          ("Graphiques :", True, BLUE_DARK),
          ("• Scatter D1 : urbanisation vs basique", False, None),
          ("• Scatter D2 : basique vs sécurisé", False, None),
          ("• Area chart évolution 2000–2017", False, None)]),
    ]

    n = len(vues)
    card_gap = 0.4
    card_w = (CONTENT_W - (n - 1) * card_gap) / n

    for i, (title, body) in enumerate(vues):
        add_card(slide, title, body,
                 MARGIN + i * (card_w + card_gap), ct, card_w, ch)

    add_page_number(slide, 10)


def build_mockup_specifiques(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "Mock-Up — Analyses spécifiques")

    vues = [
        ("Domaine 1 — Construction",
         [("3 filtres : région, pays, année", False, None),
          ("Graphique :", True, BLUE_DARK),
          ("• Scatter accès basique vs population", False, None),
          ("• Seuil < 50 % accès basique", False, None),
          ("", False, None),
          ("Cible : pays sans infrastructure", True, BLUE_MED)]),
        ("Domaine 2 — Modernisation",
         [("3 filtres : région, pays, granularité", False, None),
          ("Graphique :", True, BLUE_DARK),
          ("• Scatter basique vs sécurisé", False, None),
          ("• Seuil 50–90 % accès basique", False, None),
          ("", False, None),
          ("Cible : infra partielle à moderniser", True, BLUE_MED)]),
        ("Domaine 3 — Consulting",
         [("1 KPI stabilité · 3 filtres", False, None),
          ("Graphique :", True, BLUE_DARK),
          ("• Scatter stabilité vs accès sécurisé", False, None),
          ("", False, None),
          ("", False, None),
          ("Cible : instabilité politique forte", True, BLUE_MED)]),
    ]

    n = len(vues)
    card_gap = 0.4
    card_w = (CONTENT_W - (n - 1) * card_gap) / n

    for i, (title, body) in enumerate(vues):
        add_card(slide, title, body,
                 MARGIN + i * (card_w + card_gap), ct, card_w, ch)

    add_page_number(slide, 11)


def build_fusion(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "Fusion des données — 4 jointures Left Join")

    joins = [
        ("Population",        "Table de base\n194 pays · 2000–2017",              BLUE_DARK),
        ("+ Régions OMS",     "LEFT JOIN sur Pays\n6 régions · 194 pays",          BLUE_MED),
        ("+ Accès à l'eau",   "LEFT JOIN Pays + Année\n% basique & sécurisé",      BLUE_MED),
        ("+ Stabilité",       "LEFT JOIN Pays + Année\nWGI −3,31 à +1,76",         BLUE_MED),
        ("+ Mortalité",       "LEFT JOIN sur Pays\n2016 seulement",                GREEN),
    ]

    n = len(joins)
    arrow_w  = 0.5
    card_gap = 0.2
    card_w   = (LEFT_COL_W - (n - 1) * (arrow_w + card_gap)) / n
    card_h   = 3.8

    for i, (jname, jdesc, col) in enumerate(joins):
        left = MARGIN + i * (card_w + card_gap + arrow_w)
        _rect(slide, left, ct, card_w, card_h, col)
        _textbox(slide, jname,
                 left + 0.2, ct + 0.4, card_w - 0.4, 1.1,
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _textbox(slide, jdesc,
                 left + 0.2, ct + 1.6, card_w - 0.4, 1.9,
                 size=12, color=WHITE, align=PP_ALIGN.CENTER)
        if i < n - 1:
            _textbox(slide, "→",
                     left + card_w + card_gap / 2, ct + card_h / 2 - 0.4,
                     arrow_w, 0.8, size=20, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

    result_top = ct + card_h + 0.9
    result_lines = [
        ("dwfa_consolide.csv — table principale", True, BLUE_DARK),
        ("• 3 492 lignes · 13 colonnes · 0 doublon", False, None),
        ("• NaN documentés : accès sécurisé 50 % · mortalité 94,8 % (une seule année 2016)", False, None),
        ("", False, None),
        ("dwfa_eau_granulaire.csv — granularités Urban / Rural", True, BLUE_DARK),
        ("• 6 984 lignes · 6 colonnes", False, None),
    ]
    _multiline(slide, result_lines, MARGIN, result_top,
               LEFT_COL_W, ch - card_h - 1.2, size=18, space_after=6)

    kpi_h = (ch - 2 * 0.4) / 3
    kpis = [
        ("3 492", "Lignes consolidées",        BLUE_MED),
        ("13",    "Colonnes analytiques",       BLUE_MED),
        ("0",     "Doublon — validation OK",    GREEN),
    ]
    for i, (val, lbl, col) in enumerate(kpis):
        add_kpi_box(slide, val, lbl,
                    RIGHT_COL_L, ct + i * (kpi_h + 0.4), BOX_W, kpi_h, col)

    add_page_number(slide, 7)


def build_enrichissement(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "Enrichissement — 3 indicateurs dérivés")

    indicators = [
        ("pct_pop_rurale", "Taux de population rurale",
         "= nb_pop_rurale / nb_pop_total × 100",
         "Identifie les zones rurales non desservies, "
         "critiques pour dimensionner les interventions D1.",
         BLUE_MED),
        ("nb_deces_wash", "Décès WASH en valeur absolue",
         "= taux_mortalite_wash × nb_pop_total / 100 000",
         "Convertit le taux (pour 100 000 hab.) en valeur absolue. "
         "Comparaison directe entre pays de tailles différentes.",
         RED),
        ("pct_pop_urbaine", "Taux de population urbaine",
         "= nb_pop_urbaine / nb_pop_total × 100",
         "Indicateur clé D2 : modernisation des réseaux urbains. "
         "Disponible directement dans les données sources.",
         GREEN),
    ]

    n = len(indicators)
    card_gap = 0.4
    card_w = (CONTENT_W - (n - 1) * card_gap) / n

    for i, (col_name, title, formula, desc, col) in enumerate(indicators):
        left = MARGIN + i * (card_w + card_gap)
        _rect(slide, left, ct, card_w, ch, GRAY_BG)
        _rect(slide, left, ct, card_w, BAND_H, col)
        pad = 0.4
        title_top = ct + (BAND_H - 1.6) / 2
        txb = slide.shapes.add_textbox(Cm(left + pad), Cm(title_top),
                                       Cm(card_w - 2 * pad), Cm(1.6))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _set_run(p.add_run(), title, 16, True, WHITE)
        body_top = ct + BAND_H + 0.5
        _textbox(slide, col_name,
                 left + pad, body_top, card_w - 2 * pad, 0.9,
                 size=15, bold=True, color=col)
        _textbox(slide, formula,
                 left + pad, body_top + 1.1, card_w - 2 * pad, 0.9,
                 size=14, color=GRAY_TEXT)
        _textbox(slide, desc,
                 left + pad, body_top + 2.3, card_w - 2 * pad, ch - BAND_H - 3.0,
                 size=16, color=DARK)

    add_page_number(slide, 8)


def build_vue_mondiale(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "778 M sans accès en 2017 — l'Afrique concentre l'urgence")

    img_path = os.path.join(MOCKUP_DIR, "mockup_DWFA_vue1.png")
    slide.shapes.add_picture(img_path, Cm(MARGIN), Cm(ct), Cm(LEFT_COL_W), Cm(ch))

    bullets = [
        ("Contexte mondial 2017 :", True, BLUE_DARK),
        ("", False, None),
        ("• 89,6 % d'accès basique mondial", False, None),
        ("• Afrique : 62,4 % — région la plus critique", False, None),
        ("• +9,1 pts de progression depuis 2000", False, None),
        ("", False, None),
        ("Filtres disponibles :", True, BLUE_DARK),
        ("", False, None),
        ("• Année  ·  Région OMS  ·  Métrique", False, None),
        ("• Seuil stabilité  ·  Urbain / Rural", False, None),
    ]
    _multiline(slide, bullets, RIGHT_COL_L, ct, BOX_W, ch, size=16, space_after=6)

    add_page_number(slide, 13)


def build_vue_continentale(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "Afrique à 62,4 % — 36 points sous l'Europe")

    img_path = os.path.join(MOCKUP_DIR, "mockup_DWFA_vue2.png")
    slide.shapes.add_picture(img_path, Cm(MARGIN), Cm(ct), Cm(LEFT_COL_W), Cm(ch))

    regions = [
        ("62,4 %", "Afrique",              RED),
        ("88,9 %", "Méd. orientale",        BLUE_MED),
        ("92,5 %", "Asie du Sud-Est",        BLUE_MED),
        ("93,4 %", "Pacifique occidental",   BLUE_MED),
        ("97,6 %", "Amériques",             GREEN),
        ("98,3 %", "Europe",                GREEN),
    ]
    box_h = (ch - 5 * 0.25) / 6
    for i, (val, lbl, col) in enumerate(regions):
        top = ct + i * (box_h + 0.25)
        _rect(slide, RIGHT_COL_L, top, BOX_W, box_h, GRAY_BG)
        pad = 0.3
        _textbox(slide, val,
                 RIGHT_COL_L + pad, top + box_h * 0.1,
                 BOX_W * 0.46, box_h * 0.75,
                 size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
        _textbox(slide, lbl,
                 RIGHT_COL_L + BOX_W * 0.5, top + box_h * 0.28,
                 BOX_W * 0.46, box_h * 0.5,
                 size=13, color=GRAY_TEXT)

    add_page_number(slide, 14)


def build_vue_nationale(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "Éthiopie, RDC, Nigeria — 164 M sans accès en Afrique")

    img_path = os.path.join(MOCKUP_DIR, "mockup_DWFA_vue3.png")
    slide.shapes.add_picture(img_path, Cm(MARGIN), Cm(ct), Cm(LEFT_COL_W), Cm(ch))

    pays_data = [
        ("Éthiopie", "41,1 %", "62,7 M"),
        ("RD Congo",  "43,2 %", "46,2 M"),
        ("Nigeria",   "71,4 %", "54,6 M"),
        ("Ouganda",   "49,1 %", "21,0 M"),
        ("Tanzanie",  "56,7 %", "23,7 M"),
    ]
    box_h = (ch - 4 * 0.3) / 5
    for i, (country, pct, pop) in enumerate(pays_data):
        top = ct + i * (box_h + 0.3)
        _rect(slide, RIGHT_COL_L, top, BOX_W, box_h, GRAY_BG)
        pad = 0.3
        _textbox(slide, country,
                 RIGHT_COL_L + pad, top + 0.12, BOX_W - 2 * pad, box_h * 0.4,
                 size=15, bold=True, color=DARK)
        _textbox(slide, pct,
                 RIGHT_COL_L + pad, top + box_h * 0.52,
                 BOX_W * 0.46, box_h * 0.42,
                 size=17, bold=True, color=RED)
        _textbox(slide, pop,
                 RIGHT_COL_L + BOX_W * 0.52, top + box_h * 0.52,
                 BOX_W * 0.44, box_h * 0.42,
                 size=15, color=GRAY_TEXT)

    add_page_number(slide, 15)


def build_domaine1(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "D1 Construction — accès basique < 50 %")

    body_lines = [
        ("Critère de sélection :", True, BLUE_DARK),
        ("", False, None),
        ("Accès à l'eau basique < 50 % en 2017. Pas d'infrastructure de base — "
         "intervention de construction ab initio.", False, None),
        ("", False, None),
        ("Pays prioritaires :", True, BLUE_DARK),
        ("", False, None),
        ("• Tchad             38,7 %  →   9,2 M non couverts", False, None),
        ("• South Sudan       40,7 %  →   pop. limitée", False, None),
        ("• Éthiopie          41,1 %  →  62,7 M non couverts", False, None),
        ("• RD Congo          43,2 %  →  46,2 M non couverts", False, None),
        ("• Burkina Faso      47,9 %  →  10,0 M non couverts", False, None),
        ("", False, None),
        ("Recommandation :", True, BLUE_DARK),
        ("Éthiopie et RDC en tête par volume. Le Tchad a le taux le plus faible (38,7 %).", False, None),
    ]
    _multiline(slide, body_lines, MARGIN, ct, LEFT_COL_W, ch, size=18, space_after=6)

    kpi_h = (ch - 2 * 0.4) / 3
    kpis = [
        ("38,7 %", "Accès basique Tchad — pire taux D1",   RED),
        ("5",      "Pays avec accès < 50 % (2017)",         BLUE_MED),
        ("155 M",  "Personnes non couvertes (top 5)",       RED),
    ]
    for i, (val, lbl, col) in enumerate(kpis):
        add_kpi_box(slide, val, lbl,
                    RIGHT_COL_L, ct + i * (kpi_h + 0.4), BOX_W, kpi_h, col)

    add_page_number(slide, 17)


def build_domaine2(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "D2 Modernisation — Nigeria, Tanzanie, Kenya")

    body_lines = [
        ("Critère de sélection :", True, BLUE_DARK),
        ("", False, None),
        ("Accès basique entre 50 % et 90 % en 2017. Infrastructure existante "
         "mais incomplète — extension et sécurisation du réseau.", False, None),
        ("", False, None),
        ("Pays prioritaires :", True, BLUE_DARK),
        ("", False, None),
        ("• Nigeria      71,4 %  →  54,6 M non couverts", False, None),
        ("• Tanzanie     56,7 %  →  23,7 M non couverts", False, None),
        ("• Kenya        58,9 %  →  20,6 M non couverts", False, None),
        ("• Soudan       60,3 %  →  16,2 M non couverts", False, None),
        ("• Indonésie    89,3 %  →  28,2 M non couverts", False, None),
        ("", False, None),
        ("Recommandation :", True, BLUE_DARK),
        ("Nigeria : volume le plus élevé + écart basique/sécurisé de 51 pts.", False, None),
    ]
    _multiline(slide, body_lines, MARGIN, ct, LEFT_COL_W, ch, size=18, space_after=6)

    kpi_h = (ch - 2 * 0.4) / 3
    kpis = [
        ("71,4 %", "Accès basique Nigeria",             BLUE_MED),
        ("51 pts",  "Écart basique / sécurisé Nigeria",  RED),
        ("143 M",  "Personnes non couvertes (top 5)",   RED),
    ]
    for i, (val, lbl, col) in enumerate(kpis):
        add_kpi_box(slide, val, lbl,
                    RIGHT_COL_L, ct + i * (kpi_h + 0.4), BOX_W, kpi_h, col)

    add_page_number(slide, 18)


def build_domaine3(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "D3 Consulting — Yémen, Afghanistan, instabilité extrême")

    body_lines = [
        ("Critère de sélection :", True, BLUE_DARK),
        ("", False, None),
        ("Score WGI de stabilité politique le plus faible. Cadre institutionnel "
         "défaillant qui bloque l'amélioration de l'accès.", False, None),
        ("", False, None),
        ("Pays prioritaires :", True, BLUE_DARK),
        ("", False, None),
        ("• Yémen        score −2,94  ·  accès 63,5 %", False, None),
        ("• Afghanistan  score −2,80  ·  accès 67,1 %", False, None),
        ("• Syrie        score −2,62  ·  accès 97,2 %  (fragile)", False, None),
        ("• South Sudan  score −2,45  ·  accès 40,7 %  (D1 + D3)", False, None),
        ("• Pakistan     score −2,41  ·  accès 91,5 %", False, None),
        ("", False, None),
        ("Recommandation :", True, BLUE_DARK),
        ("Yémen : priorité absolue. South Sudan cumule urgences D1 et D3.", False, None),
    ]
    _multiline(slide, body_lines, MARGIN, ct, LEFT_COL_W, ch, size=18, space_after=6)

    kpi_h = (ch - 2 * 0.4) / 3
    kpis = [
        ("−2,94",  "Score WGI Yémen — pire stabilité",    RED),
        ("5",      "Pays avec WGI < −2,4 en 2017",        BLUE_MED),
        ("63,5 %", "Accès basique Yémen",                  RED),
    ]
    for i, (val, lbl, col) in enumerate(kpis):
        add_kpi_box(slide, val, lbl,
                    RIGHT_COL_L, ct + i * (kpi_h + 0.4), BOX_W, kpi_h, col)

    add_page_number(slide, 19)


def build_synthese(prs):
    slide = _blank_slide(prs)
    ct, ch = add_title(slide, "Synthèse — Résultats et recommandations DWFA")

    points = [
        ("Exploration & Nettoyage",
         [("• 5 sources → 1 table (3 492 lignes, 13 colonnes)", False, None),
          ("• 0 doublon — tests assert automatiques", False, None),
          ("• NaN documentés : sécurisé 50 % · mortalité 94,8 %", False, None)]),
        ("Analyse globale",
         [("• 778 M sans accès basique en 2017", False, None),
          ("• Afrique 62,4 % — 36 pts sous l'Europe (98,3 %)", False, None),
          ("• 868 000 décès WASH — 54 % en Afrique", False, None)]),
        ("Analyse des domaines",
         [("• D1 : Éthiopie 62,7 M + RDC 46,2 M non couverts", False, None),
          ("• D2 : Nigeria — 54,6 M, écart 51 pts", False, None),
          ("• D3 : Yémen −2,94 · Afghanistan −2,80", False, None)]),
        ("Outil de visualisation",
         [("• Tableau Public — 3 dashboards liés", False, None),
          ("• Filtres : région, pays, année, granularité", False, None),
          ("• Publication sans infrastructure", False, None)]),
    ]

    n = len(points)
    card_gap = 0.4
    card_w = (CONTENT_W - (n - 1) * card_gap) / n

    for i, (title, body) in enumerate(points):
        add_card(slide, title, body,
                 MARGIN + i * (card_w + card_gap), ct, card_w, ch)

    add_page_number(slide, 20)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)

    build_cover(prs)                                                    # 1
    build_objectif(prs)                                                 # 2
    build_outils(prs)                                                   # 3
    build_etapes(prs)                                                   # 4
    add_section_divider(prs, "Exploration & Nettoyage",
                        "des données", 5)                               # 5
    build_exploration_fichiers(prs, 6)                                  # 6
    build_fusion(prs)                                                   # 7

    output = os.path.join(BASE, "presentation_DWFA.pptx")
    prs.save(output)
    print(f"presentation_DWFA.pptx — {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
